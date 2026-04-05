from __future__ import annotations

import sys
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory

from openpyxl import Workbook, load_workbook
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = PROJECT_ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from excel_helpers import inspect_workbook, load_workbook_for_inspection
from generate_pptx import generate_plan_pptx


def build_workbook_with_valid_plan_and_blank_anthro(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Aaron Lopez"
    history["C5"] = "12345678"
    history["C10"] = "Masculino"

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(
        [
            "COMIDA",
            "LACTEOS",
            "VEGETALES",
            "FRUTAS",
            "ALMIDONES",
            "PROTEINAS",
            "GRASAS",
        ]
    )
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR", "CONTROL_1"])
    anthro.append(["RESUMEN", "Evaluación", None, None])
    anthro.append(["RESUMEN", "Fecha", None, None])
    anthro.append(["RESUMEN", "Peso (Kg)", None, None])
    anthro.append(["RESUMEN", "Talla parada (cm)", None, None])
    anthro.append(["RESUMEN", "% Grasa (Carter 1986)", None, None])
    anthro.append(["RESUMEN", "Kg de Grasa", None, None])
    anthro.append(["MEDIDAS", "Fecha de evaluación", None, None])
    anthro.append(["MEDIDAS", "Peso actual (kg)", None, None])
    anthro.append(["MEDIDAS", "Talla (m)", None, None])
    anthro.append(["MEDIDAS", "Talla (cm)", None, None])

    wb.save(path)


def add_guide_examples_sheet(path: Path) -> None:
    wb = load_workbook(path)
    examples = wb.create_sheet("EJEMPLOS_COMIDAS")
    examples.append(
        [
            "COMIDA",
            "LACTEOS",
            "VEGETALES",
            "FRUTAS",
            "ALMIDONES",
            "PROTEINAS",
            "GRASAS",
            "OBSERVACION",
        ]
    )
    examples.append(["PRE", None, None, "ej: fruta", "ej: almidon", "ej: proteina", None, "guia: reemplazar"])
    examples.append(["DES", None, None, "ej: fruta", "ej: almidon", "ej: proteina", "ej: grasa", "guia: reemplazar"])
    wb.save(path)


def add_observation_only_examples_sheet(path: Path) -> None:
    wb = load_workbook(path)
    examples = wb.create_sheet("EJEMPLOS_COMIDAS")
    examples.append(
        [
            "COMIDA",
            "LACTEOS",
            "VEGETALES",
            "FRUTAS",
            "ALMIDONES",
            "PROTEINAS",
            "GRASAS",
            "OBSERVACION",
        ]
    )
    examples.append(["DES", None, None, None, None, None, None, "Ajustar segun tolerancia del paciente."])
    wb.save(path)


def build_formula_based_plan_workbook(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "='DATA'!A1"
    history["C5"] = "='DATA'!A2"
    history["C7"] = "='DATA'!A3"
    history["C10"] = "='DATA'!A4"
    history["I8"] = "='DATA'!A5"

    data = wb.create_sheet("DATA")
    data["A1"] = "Aaron Lopez"
    data["A2"] = "12345678"
    data["A3"] = 31
    data["A4"] = "Masculino"
    data["A5"] = "Fitness"
    data["B1"] = 1
    data["B2"] = 0
    data["B3"] = 1
    data["B4"] = 2
    data["B5"] = 1
    data["B6"] = 1

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(
        [
            "COMIDA",
            "LACTEOS",
            "VEGETALES",
            "FRUTAS",
            "ALMIDONES",
            "PROTEINAS",
            "GRASAS",
        ]
    )
    plan.append(
        [
            "DES",
            "='DATA'!B1",
            "='DATA'!B2",
            "='DATA'!B3",
            "='DATA'!B4",
            "='DATA'!B5",
            "='DATA'!B6",
        ]
    )

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR"])
    anthro.append(["RESUMEN", "Evaluación", None])
    anthro.append(["RESUMEN", "Fecha", None])
    anthro.append(["RESUMEN", "Peso (Kg)", None])
    anthro.append(["RESUMEN", "Talla parada (cm)", None])
    anthro.append(["RESUMEN", "% Grasa (Carter 1986)", None])
    anthro.append(["RESUMEN", "Kg de Grasa", None])
    anthro.append(["MEDIDAS", "Fecha de evaluación", None])
    anthro.append(["MEDIDAS", "Peso actual (kg)", None])
    anthro.append(["MEDIDAS", "Talla (m)", None])
    anthro.append(["MEDIDAS", "Talla (cm)", None])

    wb.save(path)


class PlanGenerationRegressionTest(unittest.TestCase):
    def test_plan_generation_ignores_anthro_blockers(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Plan Alimentacion.pptx"
            build_workbook_with_valid_plan_and_blank_anthro(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            self.assertTrue(parsed_data.has_blocking_issues)

            generated_path = generate_plan_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

    def test_plan_generation_removes_dummy_example_text_when_examples_sheet_is_missing(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Plan Alimentacion.pptx"
            build_workbook_with_valid_plan_and_blank_anthro(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            generated_path = generate_plan_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            presentation = Presentation(output_path)
            slides = list(presentation.slides)[2:7]
            meal_slide_text = "\n".join(
                shape.text
                for slide in slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertNotIn("EJEMPLO:", meal_slide_text)

    def test_plan_generation_ignores_guide_examples_sheet(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Plan Alimentacion.pptx"
            build_workbook_with_valid_plan_and_blank_anthro(excel_path)
            add_guide_examples_sheet(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            generated_path = generate_plan_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            presentation = Presentation(output_path)
            slides = list(presentation.slides)[2:7]
            meal_slide_text = "\n".join(
                shape.text
                for slide in slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertNotIn("EJEMPLO:", meal_slide_text)

    def test_plan_generation_uses_note_when_example_has_only_observation(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Plan Alimentacion.pptx"
            build_workbook_with_valid_plan_and_blank_anthro(excel_path)
            add_observation_only_examples_sheet(excel_path)

            generated_path = generate_plan_pptx(
                excel_path=excel_path,
                output_path=output_path,
            )

            self.assertEqual(generated_path, output_path)
            presentation = Presentation(output_path)
            slides = list(presentation.slides)[2:7]
            meal_slide_text = "\n".join(
                shape.text
                for slide in slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            self.assertNotIn("EJEMPLO: |", meal_slide_text)
            self.assertIn("NOTA:", meal_slide_text)

    def test_formula_based_plan_inspection_and_generation_without_cached_values(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Plan Alimentacion.pptx"
            build_formula_based_plan_workbook(excel_path)

            workbook = load_workbook_for_inspection(excel_path)
            try:
                parsed_data = inspect_workbook(workbook)
            finally:
                workbook.close()

            self.assertEqual(parsed_data.patient.name, "Aaron Lopez")
            self.assertEqual(parsed_data.patient.age, "31")
            self.assertEqual(parsed_data.meal_distribution["DES"]["L"], 1.0)
            self.assertEqual(parsed_data.meal_distribution["DES"]["A"], 2.0)

            generated_path = generate_plan_pptx(
                excel_path=excel_path,
                output_path=output_path,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())


if __name__ == "__main__":
    unittest.main()
