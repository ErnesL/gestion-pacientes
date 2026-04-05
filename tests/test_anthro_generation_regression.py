from __future__ import annotations

import unittest
from datetime import datetime, time
from pathlib import Path
import sys
from tempfile import TemporaryDirectory

from openpyxl import Workbook, load_workbook
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SRC_DIR = PROJECT_ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from generate_anthro_pptx import generate_anthro_pptx
from excel_helpers import inspect_workbook, load_workbook_for_inspection


PLAN_HEADERS = [
    "COMIDA",
    "LACTEOS",
    "VEGETALES",
    "FRUTAS",
    "ALMIDONES",
    "PROTEINAS",
    "GRASAS",
]

SUMMARY_ROWS = [
    ("Evaluación", "1era evaluación"),
    ("Fecha", datetime(2026, 3, 25)),
    ("Peso (Kg)", 73.1),
    ("Talla parada (cm)", 173.5),
    ("%grasa carter", 26.0263),
    ("% Grasa (Durnin y W. 1974)", 35.56225646887326),
    ("Interpretación", "Sobrepeso II"),
    ("Kg de Masa Magra", 54.0747747),
    ("Kg de Grasa", 19.0252253),
    ("Masa Muscular (Kg)", 29.19015240045747),
    ("Masa Adiposa (Kg)", 31.7276355194793),
    ("Sumatoria de 6 pliegues", 145),
    ("Somatotipo", "Endomorfo"),
]

MEASUREMENT_ROWS = [
    ("Fecha de evaluación", datetime(2026, 3, 25)),
    ("Peso actual (kg)", 73.1),
    ("Talla (m)", 1.735),
    ("Talla (cm)", 173.5),
    ("Circunferencias(cm)", None),
    ("Brazo relajado (cm)", 31.5),
    ("Brazo Flexionado en Tensión (cm)", 32.5),
    ("Antebrazo máximo (cm)", 25),
    ("Tórax (Mesoesternal) (cm)", 93.8),
    ("Cintura mínimo (cm)", 80.7),
    ("Cadera máximo (cm)", 103.9),
    ("Muslo máximo (cm)", 63.6),
    ("Muslo medial (cm)", 55.7),
    ("Pantorrilla máximo (cm)", 35.1),
    ("Pliegues (mm)", None),
    ("Bíceps (mm)", 8.5),
    ("Tríceps (mm)", 21.5),
    ("Subescapular (mm)", 27.5),
    ("Ileo-crestal (mm)", 33.5),
    ("Supra-espinal (mm)", 22.5),
    ("Abdominal (mm)", 27.5),
    ("Muslo frontal (mm)", 32.5),
    ("Pantorrilla (mm)", 13.5),
    ("Diametros óseos (cm)", None),
    ("Humeral (Biepicondilar)", 6.1),
    ("Femoral (Biepicondilar)", 8.8),
]

SUMMARY_ROWS_MULTI = [
    ("Evaluación", ["1era evaluación", "Control 1°", "Control 2°"]),
    ("Fecha", [datetime(2026, 3, 25), datetime(2026, 4, 25), datetime(2026, 5, 25)]),
    ("Peso (Kg)", [73.1, 72.4, 71.8]),
    ("Talla parada (cm)", [173.5, 173.5, 173.5]),
    ("% Grasa (Carter 1986)", [26.0263, 24.4, 22.8]),
    ("% Grasa (Durnin y W. 1974)", [35.56225646887326, 33.2, 31.9]),
    ("Interpretación", ["Sobrepeso II", "Sobrepeso", "Saludable/Normal"]),
    ("Kg de Masa Magra", [54.0747747, 54.73, 55.43]),
    ("Kg de Grasa", [19.0252253, 17.67, 16.37]),
    ("Masa Muscular (Kg)", [29.19015240045747, 29.9, 30.4]),
    ("Masa Adiposa (Kg)", [31.7276355194793, 29.8, 28.1]),
    ("Sumatoria de 6 pliegues", [145, 136, 128]),
    ("Somatotipo", ["Endomorfo", "Meso-Endomorfo", "Mesomorfo"]),
]

MEASUREMENT_ROWS_MULTI = [
    ("Fecha de evaluación", [datetime(2026, 3, 25), datetime(2026, 4, 25), datetime(2026, 5, 25)]),
    ("Peso actual (kg)", [73.1, 72.4, 71.8]),
    ("Talla (m)", [1.735, 1.735, 1.735]),
    ("Talla (cm)", [173.5, 173.5, 173.5]),
    ("Circunferencias(cm)", [None, None, None]),
    ("Brazo relajado (cm)", [31.5, 31.2, 31.0]),
    ("Brazo Flexionado en Tensión (cm)", [32.5, 32.3, 32.1]),
    ("Antebrazo máximo (cm)", [25, 24.8, 24.6]),
    ("Tórax (Mesoesternal) (cm)", [93.8, 92.9, 92.1]),
    ("Cintura mínimo (cm)", [80.7, 78.4, 76.8]),
    ("Cadera máximo (cm)", [103.9, 102.5, 101.1]),
    ("Muslo máximo (cm)", [63.6, 62.9, 62.1]),
    ("Muslo medial (cm)", [55.7, 55.0, 54.4]),
    ("Pantorrilla máximo (cm)", [35.1, 34.9, 34.8]),
    ("Pliegues (mm)", [None, None, None]),
    ("Bíceps (mm)", [8.5, 8.0, 7.8]),
    ("Tríceps (mm)", [21.5, 20.0, 18.8]),
    ("Subescapular (mm)", [27.5, 26.0, 24.4]),
    ("Ileo-crestal (mm)", [33.5, 31.7, 29.8]),
    ("Supra-espinal (mm)", [22.5, 21.2, 19.9]),
    ("Abdominal (mm)", [27.5, 25.1, 23.0]),
    ("Muslo frontal (mm)", [32.5, 30.8, 29.1]),
    ("Pantorrilla (mm)", [13.5, 12.8, 12.0]),
    ("Diametros óseos (cm)", [None, None, None]),
    ("Humeral (Biepicondilar)", [6.1, 6.1, 6.1]),
    ("Femoral (Biepicondilar)", [8.8, 8.8, 8.8]),
]


def build_client_style_workbook(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Victoria Juliac"
    history["C5"] = "30322716"
    history["C7"] = 23
    history["C10"] = "Femenino"
    history["I8"] = ""

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR"])
    for label, value in SUMMARY_ROWS:
        anthro.append(["RESUMEN", label, value])
    for label, value in MEASUREMENT_ROWS:
        anthro.append(["MEDIDAS", label, value])

    wb.save(path)


def build_multi_consultation_workbook(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Victoria Juliac"
    history["C5"] = "30322716"
    history["C7"] = 23
    history["C10"] = "Femenino"
    history["I8"] = "Fitness"

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR", "CONTROL_1", "CONTROL_2"])
    for label, values in SUMMARY_ROWS_MULTI:
        anthro.append(["RESUMEN", label, *values])
    for label, values in MEASUREMENT_ROWS_MULTI:
        anthro.append(["MEDIDAS", label, *values])

    wb.save(path)


def build_workbook_with_missing_age_and_invalid_trailing_control(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Victoria Juliac"
    history["C5"] = "30322716"
    history["C10"] = "Femenino"
    history["I8"] = "Fitness"

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR", "CONTROL_1", "CONTROL_2"])

    summary_rows = [
        (
            label,
            [*values[:2], invalid_value],
        )
        for (label, values), invalid_value in zip(
            SUMMARY_ROWS_MULTI,
            [
                "Control 2°",
                time(0, 0),
                0,
                0,
                3.5,
                0,
                "Saludable/Normal",
                0,
                0,
                0,
                0,
                0,
                "Mesomorfo",
            ],
        )
    ]
    measurement_rows = [
        (
            label,
            [*values[:2], invalid_value],
        )
        for (label, values), invalid_value in zip(
            MEASUREMENT_ROWS_MULTI,
            [
                time(0, 0),
                0,
                0,
                0,
                None,
                0,
                0,
                0,
                0,
                0,
                0,
                0,
                0,
                0,
                None,
                0,
                0,
                0,
                0,
                0,
                0,
                0,
                0,
                None,
                0,
                0,
            ],
        )
    ]
    for label, values in summary_rows:
        anthro.append(["RESUMEN", label, *values])
    for label, values in measurement_rows:
        anthro.append(["MEDIDAS", label, *values])

    wb.save(path)


def build_workbook_with_blank_anthro(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Aaron Lopez"
    history["C5"] = "12345678"
    history["C10"] = "Masculino"

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR", "CONTROL_1"])
    anthro.append(["RESUMEN", "Evaluación", None, None])
    anthro.append(["RESUMEN", "Fecha", None, None])
    anthro.append(["RESUMEN", "Peso (Kg)", None, None])
    anthro.append(["RESUMEN", "Talla parada (cm)", None, None])
    anthro.append(["RESUMEN", "% Grasa (Carter 1986)", None, None])
    anthro.append(["RESUMEN", "Kg de Grasa", None, None])
    anthro.append(["MEDIDAS", "Fecha de evaluación", None, None])
    anthro.append(["MEDIDAS", "Peso actual (kg)", None, None])
    anthro.append(["MEDIDAS", "Talla (m)", None, None])
    anthro.append(["MEDIDAS", "Talla (cm)", None, None])

    wb.save(path)


def build_formula_based_anthro_workbook(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "='DATA'!A1"
    history["C5"] = "='DATA'!A2"
    history["C7"] = "='DATA'!A3"
    history["C10"] = "='DATA'!A4"
    history["I8"] = "='DATA'!A5"

    data = wb.create_sheet("DATA")
    data["A1"] = "Victoria Juliac"
    data["A2"] = "30322716"
    data["A3"] = 23
    data["A4"] = "Femenino"
    data["A5"] = "Fitness"
    data["B1"] = datetime(2026, 3, 25)
    data["B2"] = 73.1
    data["B3"] = 173.5
    data["B4"] = 26.0263
    data["B5"] = 19.0252253
    data["B6"] = 1.735

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR"])
    formula_summary_values = {
        "Fecha": "='DATA'!B1",
        "Peso (Kg)": "='DATA'!B2",
        "Talla parada (cm)": "='DATA'!B3",
        "% Grasa (Carter 1986)": "='DATA'!B4",
        "Kg de Grasa": "='DATA'!B5",
    }
    formula_measurement_values = {
        "Fecha de evaluación": "='DATA'!B1",
        "Peso actual (kg)": "='DATA'!B2",
        "Talla (m)": "='DATA'!B6",
        "Talla (cm)": "='DATA'!B3",
    }
    for label, value in SUMMARY_ROWS:
        anthro.append(["RESUMEN", label, formula_summary_values.get(label, value)])
    for label, value in MEASUREMENT_ROWS:
        anthro.append(["MEDIDAS", label, formula_measurement_values.get(label, value)])

    wb.save(path)


def build_workbook_with_blank_optional_rows(path: Path) -> None:
    wb = Workbook()

    history = wb.active
    history.title = "HISTORIA"
    history["C4"] = "Andrea Roo"
    history["C5"] = "12345678"
    history["C7"] = 30
    history["C10"] = "Femenino"
    history["I8"] = "Fitness"

    plan = wb.create_sheet("PLAN_ALIMENTACION_TEMPLATE")
    plan.append(PLAN_HEADERS)
    plan.append(["DES", 1, 0, 1, 2, 1, 1])

    summary_overrides = {
        "% Grasa (Durnin y W. 1974)": None,
        "Interpretación": None,
    }
    measurement_overrides = {
        "Circunferencias(cm)": None,
        "Pliegues (mm)": None,
        "Diametros óseos (cm)": None,
    }

    anthro = wb.create_sheet("ANTROPOMETRIA_TEMPLATE")
    anthro.append(["SECCION", "ETIQUETA", "VALOR"])
    for label, value in SUMMARY_ROWS:
        anthro.append(["RESUMEN", label, summary_overrides.get(label, value)])
    for label, value in MEASUREMENT_ROWS:
        anthro.append(["MEDIDAS", label, measurement_overrides.get(label, value)])

    wb.save(path)


class AnthroGenerationRegressionTest(unittest.TestCase):
    def test_generation_accepts_pct_grasa_carter_alias(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Victoria Juliac.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_client_style_workbook(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            self.assertFalse(parsed_data.has_blocking_issues)
            self.assertEqual(parsed_data.anthro_data.pct_grasa_carter, "26,03")

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

    def test_generation_builds_dynamic_tables_for_multiple_consultations(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Victoria Juliac.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_multi_consultation_workbook(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            self.assertFalse(parsed_data.has_blocking_issues)
            self.assertEqual(parsed_data.anthro_data.peso_corporal_kg, "71,80")
            self.assertEqual(parsed_data.anthro_data.pct_grasa_carter, "22,80")
            self.assertEqual(
                parsed_data.anthro_data.table_resumen[0],
                ["Evaluación", "1era evaluación", "Control 1°", "Control 2°"],
            )
            self.assertEqual(len(parsed_data.anthro_data.table_medidas[0]), 4)

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

            presentation = Presentation(output_path)
            summary_table = next(
                shape.table for shape in presentation.slides[2].shapes if shape.has_table
            )
            measurements_table = next(
                shape.table for shape in presentation.slides[3].shapes if shape.has_table
            )

            self.assertEqual(len(summary_table.columns), 4)
            self.assertEqual(len(measurements_table.columns), 4)
            self.assertEqual(
                [summary_table.cell(0, idx).text for idx in range(4)],
                ["Evaluación", "1era evaluación", "Control 1°", "Control 2°"],
            )
            self.assertEqual(
                [measurements_table.cell(0, idx).text for idx in range(4)],
                ["Fecha de evaluación", "25/3/26", "25/4/26", "25/5/26"],
            )

            slide_text = "\n".join(
                shape.text for shape in presentation.slides[2].shapes if hasattr(shape, "text")
            )
            self.assertIn("71,80", slide_text)
            self.assertIn("22,80", slide_text)

    def test_generation_ignores_invalid_trailing_controls_and_allows_missing_age(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Victoria Juliac.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_workbook_with_missing_age_and_invalid_trailing_control(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            self.assertFalse(parsed_data.has_blocking_issues)
            self.assertEqual(parsed_data.patient.age, "")
            self.assertEqual(parsed_data.anthro_data.peso_corporal_kg, "72,40")
            self.assertEqual(parsed_data.anthro_data.masa_grasa_kg, "17,67")
            self.assertEqual(
                parsed_data.anthro_data.table_resumen[0],
                ["Evaluación", "1era evaluación", "Control 1°"],
            )

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

    def test_generation_still_exports_when_anthro_has_errors(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Aaron Lopez.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_workbook_with_blank_anthro(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            self.assertTrue(parsed_data.has_blocking_issues)

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

    def test_formula_based_anthro_generation_without_cached_values(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Victoria Juliac.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_formula_based_anthro_workbook(excel_path)

            workbook = load_workbook_for_inspection(excel_path)
            try:
                parsed_data = inspect_workbook(workbook)
            finally:
                workbook.close()

            self.assertFalse(parsed_data.has_blocking_issues)
            self.assertEqual(parsed_data.patient.name, "Victoria Juliac")
            self.assertEqual(parsed_data.anthro_data.peso_corporal_kg, "73,10")
            self.assertEqual(parsed_data.anthro_data.estatura_m, "1,74")

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
            )

            self.assertEqual(generated_path, output_path)
            self.assertTrue(output_path.exists())

    def test_generation_keeps_blank_value_rows_without_inflating_table_height(self) -> None:
        with TemporaryDirectory() as tmpdir:
            temp_dir = Path(tmpdir)
            excel_path = temp_dir / "Historia Clínica - Andrea Roo.xlsx"
            output_path = temp_dir / "Informe Antropometrico.pptx"
            build_workbook_with_blank_optional_rows(excel_path)

            workbook = load_workbook(excel_path, data_only=True)
            parsed_data = inspect_workbook(workbook)

            summary_labels = [row[0] for row in parsed_data.anthro_data.table_resumen]
            measurement_labels = [row[0] for row in parsed_data.anthro_data.table_medidas]

            self.assertIn("% Grasa (Durnin y W. 1974)", summary_labels)
            self.assertIn("Interpretación", summary_labels)
            self.assertIn("Circunferencias(cm)", measurement_labels)
            self.assertIn("Pliegues (mm)", measurement_labels)
            self.assertIn("Diametros óseos (cm)", measurement_labels)
            self.assertEqual(
                parsed_data.anthro_data.table_resumen[summary_labels.index("% Grasa (Durnin y W. 1974)")][1],
                "",
            )
            self.assertEqual(
                parsed_data.anthro_data.table_resumen[summary_labels.index("Interpretación")][1],
                "",
            )
            self.assertEqual(
                parsed_data.anthro_data.table_medidas[measurement_labels.index("Circunferencias(cm)")][1],
                "",
            )
            self.assertEqual(
                parsed_data.anthro_data.table_medidas[measurement_labels.index("Pliegues (mm)")][1],
                "",
            )
            self.assertEqual(
                parsed_data.anthro_data.table_medidas[measurement_labels.index("Diametros óseos (cm)")][1],
                "",
            )

            generated_path = generate_anthro_pptx(
                excel_path=excel_path,
                output_path=output_path,
                parsed_data=parsed_data,
            )

            self.assertEqual(generated_path, output_path)
            presentation = Presentation(output_path)
            summary_table = next(
                shape.table for shape in presentation.slides[2].shapes if shape.has_table
            )
            measurements_table = next(
                shape.table for shape in presentation.slides[3].shapes if shape.has_table
            )
            template = Presentation(
                PROJECT_ROOT / "templates" / "informe-antropometrico-base.pptx"
            )
            summary_template = next(
                shape.table for shape in template.slides[2].shapes if shape.has_table
            )
            measurements_template = next(
                shape.table for shape in template.slides[3].shapes if shape.has_table
            )
            summary_table_labels = [
                summary_table.cell(row_idx, 0).text
                for row_idx in range(len(summary_table.rows))
            ]
            measurements_table_labels = [
                measurements_table.cell(row_idx, 0).text
                for row_idx in range(len(measurements_table.rows))
            ]

            self.assertIn("% Grasa (Durnin y W. 1974)", summary_table_labels)
            self.assertIn("Interpretación", summary_table_labels)
            self.assertIn("Circunferencias(cm)", measurements_table_labels)
            self.assertIn("Pliegues (mm)", measurements_table_labels)
            self.assertIn("Diametros óseos (cm)", measurements_table_labels)
            self.assertEqual(
                summary_table.cell(summary_table_labels.index("% Grasa (Durnin y W. 1974)"), 1).text,
                "",
            )
            self.assertEqual(
                summary_table.cell(summary_table_labels.index("Interpretación"), 1).text,
                "",
            )
            self.assertEqual(
                measurements_table.cell(measurements_table_labels.index("Circunferencias(cm)"), 1).text,
                "",
            )
            self.assertEqual(
                measurements_table.cell(measurements_table_labels.index("Pliegues (mm)"), 1).text,
                "",
            )
            self.assertEqual(
                measurements_table.cell(measurements_table_labels.index("Diametros óseos (cm)"), 1).text,
                "",
            )
            summary_blank_paragraph = summary_table.cell(
                summary_table_labels.index("Interpretación"), 1
            ).text_frame.paragraphs[0]
            summary_sample_paragraph = summary_template.cell(
                summary_table_labels.index("Interpretación"), 1
            ).text_frame.paragraphs[0]
            measurements_blank_paragraph = measurements_table.cell(
                measurements_table_labels.index("Circunferencias(cm)"), 1
            ).text_frame.paragraphs[0]
            measurements_sample_paragraph = measurements_template.cell(
                1, 1
            ).text_frame.paragraphs[0]
            self.assertEqual(len(summary_blank_paragraph.runs), 1)
            self.assertEqual(len(measurements_blank_paragraph.runs), 1)
            self.assertEqual(
                summary_blank_paragraph.runs[0].font.size,
                summary_sample_paragraph.runs[0].font.size,
            )
            self.assertEqual(
                measurements_blank_paragraph.runs[0].font.size,
                measurements_sample_paragraph.runs[0].font.size,
            )
            self.assertIn("endParaRPr", summary_table.cell(
                summary_table_labels.index("Interpretación"), 1
            )._tc.xml)
            self.assertIn("kern=\"1200\"", summary_table.cell(
                summary_table_labels.index("Interpretación"), 1
            )._tc.xml)
            self.assertIn("endParaRPr", measurements_table.cell(
                measurements_table_labels.index("Circunferencias(cm)"), 1
            )._tc.xml)
            self.assertIn("kern=\"1200\"", measurements_table.cell(
                measurements_table_labels.index("Circunferencias(cm)"), 1
            )._tc.xml)
            self.assertEqual(
                summary_table.rows[len(summary_table.rows) - 1].height,
                summary_template.rows[len(summary_table.rows) - 1].height,
            )
            self.assertEqual(
                measurements_table.rows[len(measurements_table.rows) - 1].height,
                measurements_template.rows[len(measurements_table.rows) - 1].height,
            )


if __name__ == "__main__":
    unittest.main()
