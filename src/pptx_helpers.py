from __future__ import annotations

import re
from copy import deepcopy
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple

from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class TableColumnMap:
    shape: object
    key_to_index: Dict[str, int]
    removed_cols: List[int]
    group_to_index: Dict[str, int]
    image_slots: Dict[str, Tuple[int, int, int, int]]


COL_MARKER_RE = re.compile(r"{{COL[:_](?P<key>[A-Z_]+)}}")
IMG_MARKER_RE = re.compile(r"{{IMG[:_](?P<key>[A-Z_]+)}}")
SHAPE_COL_MARKER_RE = re.compile(r"COL_(?P<key>[A-Z_]+)")
SHAPE_IMG_MARKER_RE = re.compile(r"IMG_(?P<key>[A-Z_]+)")
SHAPE_KEY_SUFFIXES = ("_ARROW", "_FLECHA", "_ICON")
STACK_MARKER_RE = re.compile(
    r"STACK_(?P<group>\\d+)(?:_(?P<item>\\d+))?(?:_GAP(?P<gap>\\d+))?",
    re.IGNORECASE,
)
STACK_GAP = 120_000
EXAMPLE_PREFIX_RE = re.compile(
    r"^(?P<prefix>(?:EJEMPLO|NOTA):\s*)(?P<body>.*)$",
    re.IGNORECASE,
)

GROUP_LABELS = {
    "LACTEOS": "LACTEOS",
    "VEGETALES": "VEGETALES",
    "FRUTAS": "FRUTAS",
    "ALMIDONES": "ALMIDONES",
    "PROTEINAS": "PROTEINAS",
    "PROTEICOS": "PROTEINAS",
    "GRASAS": "GRASAS",
}


def normalize_label(text: str) -> str:
    normalized = text.upper()
    for src, dst in (
        ("Á", "A"),
        ("É", "E"),
        ("Í", "I"),
        ("Ó", "O"),
        ("Ú", "U"),
        ("Ü", "U"),
    ):
        normalized = normalized.replace(src, dst)
    return " ".join(normalized.split())


def normalize_shape_key(key: str) -> str:
    for suffix in SHAPE_KEY_SUFFIXES:
        if key.endswith(suffix):
            return key[: -len(suffix)]
    return key


def key_to_group(key: str) -> str:
    if "_" in key:
        return key.split("_", 1)[1]
    return key


def replace_in_text(text: str, replacements: Dict[str, str]) -> str:
    updated = text
    for key, value in replacements.items():
        if key in updated:
            updated = updated.replace(key, value)
    return updated


def replace_in_text_frame(text_frame, replacements: Dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue
        original = "".join(run.text for run in paragraph.runs)
        if not original:
            continue
        updated = replace_in_text(original, replacements)
        if updated == original:
            continue
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""


def set_text_frame_text(text_frame, text: str) -> None:
    if not text_frame.paragraphs:
        paragraph = text_frame.add_paragraph()
        paragraph.text = text
        return

    first_paragraph = text_frame.paragraphs[0]
    if set_prefixed_text_frame_text(first_paragraph, text):
        for paragraph in text_frame.paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
        remove_empty_paragraphs(text_frame)
        return

    if first_paragraph.runs:
        first_paragraph.runs[0].text = text
        for run in first_paragraph.runs[1:]:
            run.text = ""
    else:
        first_paragraph.text = text

    for paragraph in text_frame.paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""
    remove_empty_paragraphs(text_frame)


def set_prefixed_text_frame_text(paragraph, text: str) -> bool:
    match = EXAMPLE_PREFIX_RE.match(text)
    if match is None or len(paragraph.runs) < 2:
        return False

    # Keep the template's distinct styles for the prefix and example body.
    paragraph.runs[0].text = match.group("prefix")
    paragraph.runs[1].text = match.group("body")
    for run in paragraph.runs[2:]:
        run.text = ""
    return True


def remove_empty_paragraphs(text_frame) -> None:
    to_remove = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            to_remove.append(paragraph._element)
    for element in to_remove:
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)


def strip_col_markers(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue
        original = "".join(run.text for run in paragraph.runs)
        if not original:
            continue
        updated = COL_MARKER_RE.sub("", original)
        updated = " ".join(updated.split())
        if updated == original:
            continue
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
    remove_empty_paragraphs(text_frame)


def strip_img_markers(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            continue
        original = "".join(run.text for run in paragraph.runs)
        if not original:
            continue
        updated = IMG_MARKER_RE.sub("", original)
        updated = " ".join(updated.split())
        if updated == original:
            continue
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
    remove_empty_paragraphs(text_frame)


def should_hide_shape(shape, placeholder_values: Dict[str, int]) -> bool:
    name = getattr(shape, "name", "") or ""
    for key in SHAPE_COL_MARKER_RE.findall(name) + SHAPE_IMG_MARKER_RE.findall(name):
        normalized = normalize_shape_key(key)
        placeholder = f"{{{{{normalized}}}}}"
        if placeholder_values.get(placeholder, 0) == 0:
            return True
    return False


def remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def copy_table_style(source_table, target_table) -> None:
    source_tbl_pr = source_table._tbl.tblPr
    target_tbl_pr = target_table._tbl.tblPr

    for attr_name in list(target_tbl_pr.attrib):
        del target_tbl_pr.attrib[attr_name]
    for attr_name, attr_value in source_tbl_pr.attrib.items():
        target_tbl_pr.set(attr_name, attr_value)

    for child in list(target_tbl_pr):
        target_tbl_pr.remove(child)
    for child in source_tbl_pr:
        target_tbl_pr.append(deepcopy(child))


def compute_dynamic_table_col_widths(source_table, target_col_count: int, total_width: int) -> List[int]:
    source_widths = [col.width for col in source_table.columns]
    if target_col_count <= len(source_widths):
        widths = source_widths[:target_col_count]
        if widths:
            widths[-1] += total_width - sum(widths)
        return widths

    label_width = source_widths[0] if source_widths else int(total_width / target_col_count)
    label_width = min(label_width, int(total_width * 0.4))
    label_width = max(label_width, int(total_width * 0.28))

    value_col_count = max(target_col_count - 1, 1)
    remaining_width = max(total_width - label_width, value_col_count)
    value_width = int(remaining_width / value_col_count)
    widths = [label_width] + [value_width] * value_col_count
    widths[-1] += total_width - sum(widths)
    return widths


def compute_dynamic_table_row_heights(source_table, target_row_count: int, total_height: int) -> List[int]:
    source_heights = [row.height for row in source_table.rows]
    if not source_heights:
        height = int(total_height / max(target_row_count, 1))
        return [height] * target_row_count

    if target_row_count <= len(source_heights):
        return source_heights[:target_row_count]

    heights = source_heights + [source_heights[-1]] * (target_row_count - len(source_heights))
    current_total = sum(heights)
    if current_total <= total_height:
        heights[-1] += total_height - current_total
        return heights

    scale = total_height / current_total
    scaled_heights = [max(1, int(height * scale)) for height in heights]
    scaled_heights[-1] += total_height - sum(scaled_heights)
    return scaled_heights


def copy_paragraph_style(source_paragraph, target_paragraph) -> None:
    target_paragraph.alignment = source_paragraph.alignment
    target_paragraph.level = source_paragraph.level
    source_p = source_paragraph._p
    target_p = target_paragraph._p

    def replace_child(child_name: str) -> None:
        existing = next(
            (child for child in target_p if child.tag.rsplit("}", 1)[-1] == child_name),
            None,
        )
        if existing is not None:
            target_p.remove(existing)
        source_child = next(
            (child for child in source_p if child.tag.rsplit("}", 1)[-1] == child_name),
            None,
        )
        if source_child is None:
            return
        if child_name == "pPr":
            target_p.insert(0, deepcopy(source_child))
        else:
            target_p.append(deepcopy(source_child))

    replace_child("pPr")
    replace_child("endParaRPr")


def copy_run_style(source_run, target_run) -> None:
    source_r = source_run._r
    target_r = target_run._r
    existing = next(
        (child for child in target_r if child.tag.rsplit("}", 1)[-1] == "rPr"),
        None,
    )
    if existing is not None:
        target_r.remove(existing)
    source_rpr = next(
        (child for child in source_r if child.tag.rsplit("}", 1)[-1] == "rPr"),
        None,
    )
    if source_rpr is not None:
        target_r.insert(0, deepcopy(source_rpr))

    source_font = source_run.font
    target_font = target_run.font

    if source_font.name is not None:
        target_font.name = source_font.name
    if source_font.size is not None:
        target_font.size = source_font.size
    if source_font.bold is not None:
        target_font.bold = source_font.bold
    if source_font.italic is not None:
        target_font.italic = source_font.italic
    if source_font.underline is not None:
        target_font.underline = source_font.underline

    try:
        if source_font.color.rgb is not None:
            target_font.color.rgb = source_font.color.rgb
    except Exception:
        pass


def set_table_cell_text_from_sample(target_cell, text: str, sample_cell) -> None:
    target_cell.text = text
    target_cell.margin_left = sample_cell.margin_left
    target_cell.margin_right = sample_cell.margin_right
    target_cell.margin_top = sample_cell.margin_top
    target_cell.margin_bottom = sample_cell.margin_bottom
    target_cell.vertical_anchor = sample_cell.vertical_anchor

    target_text_frame = target_cell.text_frame
    source_text_frame = sample_cell.text_frame
    target_text_frame.word_wrap = source_text_frame.word_wrap

    if not source_text_frame.paragraphs or not target_text_frame.paragraphs:
        return

    source_paragraph = source_text_frame.paragraphs[0]
    target_paragraph = target_text_frame.paragraphs[0]
    copy_paragraph_style(source_paragraph, target_paragraph)

    if not source_paragraph.runs:
        return

    if target_paragraph.runs:
        target_run = target_paragraph.runs[0]
        target_run.text = text
        for extra_run in target_paragraph.runs[1:]:
            extra_run.text = ""
    else:
        target_run = target_paragraph.add_run()
        target_run.text = text

    copy_run_style(source_paragraph.runs[0], target_run)


def cell_has_styled_run(cell) -> bool:
    paragraphs = cell.text_frame.paragraphs
    return bool(paragraphs and paragraphs[0].runs)


def resolve_table_sample_cell(source_table, row_idx: int, col_idx: int):
    sample_row_idx = min(row_idx, len(source_table.rows) - 1)
    sample_col_idx = 0 if col_idx == 0 else min(1, len(source_table.columns) - 1)
    sample_cell = source_table.cell(sample_row_idx, sample_col_idx)
    if cell_has_styled_run(sample_cell):
        return sample_cell

    if sample_col_idx > 0:
        fallback_row_idx = 0 if row_idx == 0 else min(1, len(source_table.rows) - 1)
        fallback_cell = source_table.cell(fallback_row_idx, sample_col_idx)
        if cell_has_styled_run(fallback_cell):
            return fallback_cell

    for fallback_row_idx in range(len(source_table.rows)):
        fallback_cell = source_table.cell(fallback_row_idx, sample_col_idx)
        if cell_has_styled_run(fallback_cell):
            return fallback_cell

    return sample_cell


def replace_table_shape_with_data(slide, shape, data_rows: List[List[str]]) -> object:
    if not data_rows:
        remove_shape(shape)
        return None

    row_count = len(data_rows)
    col_count = max(len(row) for row in data_rows)
    source_table = shape.table

    new_shape = slide.shapes.add_table(
        row_count,
        col_count,
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )
    target_table = new_shape.table
    copy_table_style(source_table, target_table)

    col_widths = compute_dynamic_table_col_widths(source_table, col_count, shape.width)
    for idx, width in enumerate(col_widths):
        target_table.columns[idx].width = width

    row_heights = compute_dynamic_table_row_heights(source_table, row_count, shape.height)
    for idx, height in enumerate(row_heights):
        target_table.rows[idx].height = height
    rendered_height = sum(row_heights)
    if rendered_height < shape.height:
        new_shape.height = rendered_height

    for row_idx in range(row_count):
        for col_idx in range(col_count):
            sample_cell = resolve_table_sample_cell(source_table, row_idx, col_idx)
            value = data_rows[row_idx][col_idx] if col_idx < len(data_rows[row_idx]) else ""
            set_table_cell_text_from_sample(
                target_table.cell(row_idx, col_idx),
                value,
                sample_cell,
            )

    remove_shape(shape)
    return new_shape


def remove_table_columns(table, col_indices: List[int]) -> None:
    tbl = table._tbl
    grid = tbl.tblGrid
    for col_idx in col_indices:
        grid_cols = getattr(grid, "gridCol_lst", None)
        if grid_cols is None:
            grid_cols = list(grid.iterchildren())
        if col_idx < 0 or col_idx >= len(grid_cols):
            continue
        grid.remove(grid_cols[col_idx])
        for row in table.rows:
            cells = getattr(row._tr, "tc_lst", None)
            if cells is None:
                cells = list(row._tr.iterchildren())
            if col_idx < len(cells):
                row._tr.remove(cells[col_idx])


def remove_empty_table_rows(table, shape) -> None:
    tbl = table._tbl
    rows = list(table.rows)
    rows_to_remove = []
    for idx, row in enumerate(rows):
        if idx == 0:
            continue
        if all(not cell.text.strip() for cell in row.cells):
            rows_to_remove.append(idx)
    if not rows_to_remove:
        return
    tr_list = getattr(tbl, "tr_lst", None)
    if tr_list is None:
        tr_list = list(tbl.iterchildren())
    for idx in sorted(rows_to_remove, reverse=True):
        if idx < len(tr_list):
            tbl.remove(tr_list[idx])
    shape.height = sum(row.height for row in table.rows)


def find_image_slots(table, shape) -> Dict[str, Tuple[int, int, int, int]]:
    slots: Dict[str, Tuple[int, int, int, int]] = {}
    if not table.rows or not table.columns:
        return slots

    widths = [col.width for col in table.columns]
    heights = [row.height for row in table.rows]
    if any(h is None for h in heights):
        avg_height = int(shape.height / len(heights))
        heights = [h if h is not None else avg_height for h in heights]

    col_lefts = []
    acc = shape.left
    for width in widths:
        col_lefts.append(acc)
        acc += width

    row_tops = []
    acc = shape.top
    for height in heights:
        row_tops.append(acc)
        acc += height

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            text = cell.text
            markers = IMG_MARKER_RE.findall(text)
            if not markers:
                continue
            for key in markers:
                left = col_lefts[col_idx]
                top = row_tops[row_idx]
                slots[key] = (left, top, widths[col_idx], heights[row_idx])
            strip_img_markers(cell.text_frame)

    return slots


def replace_in_shape(
    shape,
    replacements: Dict[str, str],
    placeholder_values: Dict[str, float],
    slide_width: int,
    slide_shapes,
    table_maps: List[TableColumnMap],
) -> None:
    if should_hide_shape(shape, placeholder_values):
        remove_shape(shape)
        return

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            replace_in_shape(
                subshape,
                replacements,
                placeholder_values,
                slide_width,
                slide_shapes,
                table_maps,
            )
        return

    if shape.has_table:
        table = shape.table
        cols_to_hide = set()
        key_to_index: Dict[str, int] = {}
        group_to_index: Dict[str, int] = {}
        for col_idx in range(len(table.columns)):
            for row in table.rows:
                cell_text = row.cells[col_idx].text
                markers = COL_MARKER_RE.findall(cell_text)
                if markers:
                    for key in markers:
                        key_to_index.setdefault(key, col_idx)
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder_values.get(placeholder, 0) == 0:
                            cols_to_hide.add(col_idx)
                    strip_col_markers(row.cells[col_idx].text_frame)
                for placeholder, value in placeholder_values.items():
                    if placeholder in cell_text:
                        key_to_index.setdefault(
                            placeholder.strip("{}"), col_idx)
                        if value == 0:
                            cols_to_hide.add(col_idx)
            if table.rows:
                header_text = table.rows[0].cells[col_idx].text
                normalized_header = normalize_label(header_text)
                for label, group in GROUP_LABELS.items():
                    if label and label in normalized_header:
                        group_to_index.setdefault(group, col_idx)

        if cols_to_hide:
            remove_table_columns(table, sorted(cols_to_hide, reverse=True))
            if slide_width and shape._parent is slide_shapes:
                table_width = sum(col.width for col in table.columns)
                shape.width = table_width
                shape.left = int((slide_width - shape.width) / 2)
            group_to_index = {
                group: idx -
                sum(1 for removed in cols_to_hide if removed < idx)
                for group, idx in group_to_index.items()
            }

        remove_empty_table_rows(table, shape)
        image_slots = find_image_slots(table, shape)
        table_maps.append(
            TableColumnMap(
                shape=shape,
                key_to_index=key_to_index,
                removed_cols=sorted(cols_to_hide),
                group_to_index=group_to_index,
                image_slots=image_slots,
            )
        )

        for row in table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame, replacements)
        return

    if shape.has_text_frame:
        replace_in_text_frame(shape.text_frame, replacements)


def text_frame_contains(text_frame, tokens: List[str]) -> bool:
    for paragraph in text_frame.paragraphs:
        if any(token in paragraph.text for token in tokens):
            return True
    return False


def shape_contains_tokens(shape, tokens: List[str]) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return any(shape_contains_tokens(subshape, tokens) for subshape in shape.shapes)
    if shape.has_text_frame and text_frame_contains(shape.text_frame, tokens):
        return True
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if text_frame_contains(cell.text_frame, tokens):
                    return True
    return False


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in iter_shapes(shape.shapes):
                yield subshape


def find_shape_by_name(slide, shape_name: str) -> Optional[object]:
    for shape in iter_shapes(slide.shapes):
        if (getattr(shape, "name", "") or "") == shape_name:
            return shape
    return None


def replace_meal_example_text(slide, example_text: str) -> bool:
    replaced = False
    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        full_text = "\n".join(
            paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text
        ).strip()
        normalized_text = normalize_label(full_text)
        if "{{" in full_text and "EJEMPLO" in normalized_text:
            set_text_frame_text(shape.text_frame, example_text)
            replaced = True
            continue
        if normalized_text.startswith("EJEMPLO:"):
            set_text_frame_text(shape.text_frame, example_text)
            replaced = True
    return replaced


def align_marked_shapes(
    slide, placeholder_values: Dict[str, float], table_maps: List[TableColumnMap]
) -> None:
    if not table_maps:
        return

    table_centers: List[Tuple[TableColumnMap, Dict[str, int]]] = []
    for table_map in table_maps:
        table = table_map.shape.table
        widths = [col.width for col in table.columns]
        removed = table_map.removed_cols
        centers: Dict[str, int] = {}
        for key, original_idx in table_map.key_to_index.items():
            placeholder = f"{{{{{key}}}}}"
            if placeholder_values.get(placeholder, 0) == 0:
                continue
            shift = sum(1 for idx in removed if idx < original_idx)
            new_idx = original_idx - shift
            if new_idx < 0 or new_idx >= len(widths):
                continue
            left = table_map.shape.left + sum(widths[:new_idx])
            center = left + int(widths[new_idx] / 2)
            centers[key] = center
        for group, idx in table_map.group_to_index.items():
            if idx < 0 or idx >= len(widths):
                continue
            left = table_map.shape.left + sum(widths[:idx])
            centers[group] = left + int(widths[idx] / 2)
        table_centers.append((table_map, centers))

    if not table_centers:
        return

    for shape in iter_shapes(slide.shapes):
        name = getattr(shape, "name", "") or ""
        img_keys = SHAPE_IMG_MARKER_RE.findall(name)
        if img_keys:
            key = normalize_shape_key(img_keys[0])
            shape_center_y = shape.top + int(shape.height / 2)
            best_slot = None
            best_distance = None
            for table_map in table_maps:
                slot = table_map.image_slots.get(key)
                if slot is None:
                    slot = table_map.image_slots.get(key_to_group(key))
                if slot is None:
                    continue
                _, slot_top, _, slot_height = slot
                slot_center_y = slot_top + int(slot_height / 2)
                distance = abs(shape_center_y - slot_center_y)
                if best_distance is None or distance < best_distance:
                    best_distance = distance
                    best_slot = slot
            if best_slot is not None:
                left, top, width, height = best_slot
                shape.left = int(left + (width - shape.width) / 2)
                shape.top = int(top + (height - shape.height) / 2)
            continue
        keys = SHAPE_COL_MARKER_RE.findall(name)
        if not keys:
            continue
        key = normalize_shape_key(keys[0])
        shape_center_y = shape.top + int(shape.height / 2)
        best_center = None
        best_distance = None
        for table_map, centers in table_centers:
            center = centers.get(key)
            if center is None:
                center = centers.get(key_to_group(key))
            if center is None:
                continue
            table_center_y = table_map.shape.top + \
                int(table_map.shape.height / 2)
            distance = abs(shape_center_y - table_center_y)
            if best_distance is None or distance < best_distance:
                best_distance = distance
                best_center = center
        if best_center is None:
            continue
        shape.left = int(best_center - shape.width / 2)


def apply_vertical_stack(slide) -> None:
    groups = {}
    group_gaps = {}
    for shape in slide.shapes:
        name = getattr(shape, "name", "") or ""
        match = STACK_MARKER_RE.search(name)
        if not match:
            continue
        group = int(match.group("group"))
        item = match.group("item")
        order = int(item) if item else 0
        gap = match.group("gap")
        if gap:
            group_gaps[group] = int(gap)
        groups.setdefault(group, []).append((order, shape))

    if len(groups) < 2:
        return

    for group, items in groups.items():
        items.sort(key=lambda item: (item[0], getattr(item[1], "name", "")))
        groups[group] = [shape for _, shape in items]

    ordered_groups = sorted(groups.items(), key=lambda item: item[0])
    first_group = ordered_groups[0][1]
    first_top, first_bottom = group_visual_bounds(first_group)
    current_top = first_bottom + \
        group_gaps.get(ordered_groups[0][0], STACK_GAP)

    for group_id, shapes in ordered_groups[1:]:
        group_top, group_bottom = group_visual_bounds(shapes)
        delta = current_top - group_top
        for shape in shapes:
            shape.top = int(shape.top + delta)
        current_top = group_bottom + delta + \
            group_gaps.get(group_id, STACK_GAP)


def shape_visual_bounds(shape) -> Tuple[int, int]:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        tops = []
        bottoms = []
        for subshape in shape.shapes:
            top, bottom = shape_visual_bounds(subshape)
            tops.append(top)
            bottoms.append(bottom)
        if tops and bottoms:
            return min(tops), max(bottoms)
    return shape.top, shape.top + shape.height


def group_visual_bounds(shapes) -> Tuple[int, int]:
    tops = []
    bottoms = []
    for shape in shapes:
        top, bottom = shape_visual_bounds(shape)
        tops.append(top)
        bottoms.append(bottom)
    return min(tops), max(bottoms)


def slide_contains_tokens(slide, tokens: List[str]) -> bool:
    return any(shape_contains_tokens(shape, tokens) for shape in slide.shapes)
