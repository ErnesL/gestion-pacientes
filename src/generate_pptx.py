from __future__ import annotations

import argparse
from pathlib import Path
from typing import Dict, List

from pptx import Presentation

from excel_helpers import (
    MEAL_DEFS,
    ParsedWorkbookData,
    ValidationError,
    build_validation_warning_message,
    build_meal_replacements,
    build_replacements,
    build_totals_replacements,
    inspect_workbook,
    load_workbook_for_inspection,
)
from pptx_helpers import (
    align_marked_shapes,
    apply_vertical_stack,
    replace_meal_example_text,
    replace_in_shape,
    slide_contains_tokens,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE_PATH = PROJECT_ROOT / \
    "templates" / "plan-de-alimentacion-base.pptx"
DEFAULT_OUTPUT_PATH = PROJECT_ROOT / "output" / "Plan Alimentacion.pptx"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Genera PPTX desde Excel")
    parser.add_argument("excel", help="Ruta al archivo Excel")
    parser.add_argument(
        "--template",
        default=str(DEFAULT_TEMPLATE_PATH),
        help="Ruta al PPTX base con placeholders",
    )
    parser.add_argument(
        "--output",
        default=str(DEFAULT_OUTPUT_PATH),
        help="Ruta de salida PPTX",
    )
    return parser.parse_args()


def remove_slides_by_index(presentation: Presentation, indices: List[int]) -> None:
    if not indices:
        return
    slide_id_list = presentation.slides._sldIdLst
    for index in sorted(indices, reverse=True):
        slide_id_list.remove(slide_id_list[index])


def generate_plan_pptx(
    excel_path: Path | str,
    template_path: Path | str = DEFAULT_TEMPLATE_PATH,
    output_path: Path | str = DEFAULT_OUTPUT_PATH,
    parsed_data: ParsedWorkbookData | None = None,
) -> Path:
    excel_path = Path(excel_path)
    template_path = Path(template_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if not excel_path.exists():
        raise FileNotFoundError(f"No existe el archivo: {excel_path}")
    if not template_path.exists():
        raise FileNotFoundError(f"No existe el PPTX base: {template_path}")

    if parsed_data is None:
        wb = load_workbook_for_inspection(excel_path)
        try:
            parsed_data = inspect_workbook(wb)
        finally:
            wb.close()

    patient = parsed_data.patient
    meal_distribution = parsed_data.meal_distribution

    replacements = build_replacements(patient)
    replacements.update(build_totals_replacements(meal_distribution))
    meal_example_texts = parsed_data.meal_examples

    presentation = Presentation(str(template_path))
    slides_to_remove: List[int] = []
    meal_tokens_to_remove: List[List[str]] = []
    meal_tokens_by_name: Dict[str, List[str]] = {}
    placeholder_values: Dict[str, float] = {}

    for meal_def in MEAL_DEFS:
        meal_repl, _, include_meal, tokens, meal_placeholder_values = build_meal_replacements(
            meal_distribution, meal_def
        )
        replacements.update(meal_repl)
        example_placeholder = "{{" + meal_def["name"] + "_EJEMPLO}}"
        replacements[example_placeholder] = meal_example_texts.get(
            meal_def["name"],
            "",
        )
        placeholder_values.update(meal_placeholder_values)
        meal_tokens_by_name[meal_def["name"]] = tokens
        if not include_meal:
            meal_tokens_to_remove.append(tokens)

    for idx, slide in enumerate(presentation.slides):
        slide_meal_name = None
        for meal_name, tokens in meal_tokens_by_name.items():
            if slide_contains_tokens(slide, tokens):
                slide_meal_name = meal_name
                break
        for tokens in meal_tokens_to_remove:
            if slide_contains_tokens(slide, tokens):
                slides_to_remove.append(idx)
                break
        if idx in slides_to_remove:
            continue
        table_maps = []
        for shape in slide.shapes:
            replace_in_shape(
                shape,
                replacements,
                placeholder_values,
                presentation.slide_width,
                slide.shapes,
                table_maps,
            )
        align_marked_shapes(slide, placeholder_values, table_maps)
        apply_vertical_stack(slide)
        if slide_meal_name:
            replace_meal_example_text(
                slide,
                meal_example_texts.get(slide_meal_name, ""),
            )

    remove_slides_by_index(presentation, slides_to_remove)

    presentation.save(str(output_path))
    return output_path


def main() -> int:
    args = parse_args()
    try:
        wb = load_workbook_for_inspection(args.excel)
        try:
            parsed_data = inspect_workbook(wb)
        finally:
            wb.close()
        output_path = generate_plan_pptx(
            excel_path=args.excel,
            template_path=args.template,
            output_path=args.output,
            parsed_data=parsed_data,
        )
    except (FileNotFoundError, ValidationError) as exc:
        print(f"Error: {exc}")
        return 1
    print(f"PPTX generado: {output_path}")
    if parsed_data.issues:
        print("")
        print(build_validation_warning_message(parsed_data.issues))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
