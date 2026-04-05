from __future__ import annotations

import argparse
from datetime import date, datetime
from pathlib import Path

from pptx import Presentation

from excel_helpers import (
    ParsedWorkbookData,
    ValidationError,
    build_anthropometric_replacements,
    build_validation_warning_message,
    inspect_workbook,
    load_workbook_for_inspection,
)
from pptx_helpers import replace_in_shape, replace_table_shape_with_data


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE_PATH = PROJECT_ROOT / \
    "templates" / "informe-antropometrico-base.pptx"
DEFAULT_OUTPUT_PATH = PROJECT_ROOT / "output" / "Informe Antropometrico.pptx"


def parse_iso_date(value: str) -> date:
    try:
        return datetime.strptime(value, "%Y-%m-%d").date()
    except ValueError as exc:
        raise argparse.ArgumentTypeError(
            "Formato invalido para --today. Usa YYYY-MM-DD."
        ) from exc


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Genera Informe Antropometrico PPTX desde Excel"
    )
    parser.add_argument("excel", help="Ruta al archivo Excel")
    parser.add_argument(
        "--template",
        default=str(DEFAULT_TEMPLATE_PATH),
        help="Ruta al PPTX antropometrico base",
    )
    parser.add_argument(
        "--output",
        default=str(DEFAULT_OUTPUT_PATH),
        help="Ruta de salida PPTX",
    )
    parser.add_argument(
        "--today",
        type=parse_iso_date,
        default=None,
        help="Fecha base para calculos (YYYY-MM-DD). Si no se indica, usa hoy.",
    )
    return parser.parse_args()


def generate_anthro_pptx(
    excel_path: Path | str,
    template_path: Path | str = DEFAULT_TEMPLATE_PATH,
    output_path: Path | str = DEFAULT_OUTPUT_PATH,
    today: date | None = None,
    parsed_data: ParsedWorkbookData | None = None,
) -> Path:
    excel_path = Path(excel_path)
    template_path = Path(template_path)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if not excel_path.exists():
        raise FileNotFoundError(f"No existe el archivo: {excel_path}")
    if not template_path.exists():
        raise FileNotFoundError(f"No existe el PPTX base: {template_path}")

    if parsed_data is None:
        wb = load_workbook_for_inspection(excel_path)
        try:
            parsed_data = inspect_workbook(wb)
        finally:
            wb.close()
    anthro_data = parsed_data.anthro_data

    reference_date = today if today is not None else date.today()
    base_replacements = build_anthropometric_replacements(
        anthro_data, reference_date)

    presentation = Presentation(str(template_path))
    if len(presentation.slides) < 4:
        raise ValidationError(
            "El template antropometrico debe tener al menos 4 diapositivas."
        )

    for slide_idx, slide in enumerate(presentation.slides):
        slide_replacements = dict(base_replacements)

        table_maps = []
        for shape in list(slide.shapes):
            if slide_idx == 2 and shape.has_table:
                replace_table_shape_with_data(
                    slide,
                    shape,
                    anthro_data.table_resumen,
                )
                continue
            if slide_idx == 3 and shape.has_table:
                replace_table_shape_with_data(
                    slide,
                    shape,
                    anthro_data.table_medidas,
                )
                continue
            replace_in_shape(
                shape,
                slide_replacements,
                {},
                presentation.slide_width,
                slide.shapes,
                table_maps,
            )

    presentation.save(str(output_path))
    return output_path


def main() -> int:
    args = parse_args()
    try:
        wb = load_workbook_for_inspection(args.excel)
        try:
            parsed_data = inspect_workbook(wb)
        finally:
            wb.close()
        output_path = generate_anthro_pptx(
            excel_path=args.excel,
            template_path=args.template,
            output_path=args.output,
            today=args.today,
            parsed_data=parsed_data,
        )
    except (FileNotFoundError, ValidationError) as exc:
        print(f"Error: {exc}")
        return 1
    print(f"PPTX antropometrico generado: {output_path}")
    if parsed_data.issues:
        print("")
        print(build_validation_warning_message(parsed_data.issues))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
